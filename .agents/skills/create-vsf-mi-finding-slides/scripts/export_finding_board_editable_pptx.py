#!/usr/bin/env python3
"""Export an approved WR3 Finding Board as a fully editable native PPTX."""

from __future__ import annotations

import argparse
import importlib.util
import re
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


W, H = 13.333333, 7.5
RED = "E5002B"
DEEP_RED = "8F151A"
TEXT = "202124"
MUTED = "666666"
BORDER = "DDDCD8"
SOFT = "F8F8F7"
SIGNAL_FILL = "EAF3FF"
SIGNAL_LINE = "BED5EF"
OPP_FILL = "DDF5EE"
OPP_LINE = "91D9C3"
THREAT_FILL = "FDE8E8"
THREAT_LINE = "F4B4B4"
BLUE_TEXT = "123B66"
FONT_HEAD = "VSF Pro"
FONT_BODY = "Lexend"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def load_html_generator(path: Path):
    spec = importlib.util.spec_from_file_location("finding_html", path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Cannot load generator: {path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def add_shape(slide, x, y, w, h, fill="FFFFFF", line=BORDER, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.75)
    return shape


def set_text_frame(tf, margin=0.06, valign=MSO_ANCHOR.TOP, fit=True):
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return tf


def style_run(run, size=12, color=TEXT, bold=False, font=FONT_BODY, italic=False, underline=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    return run


def text_box(slide, x, y, w, h, text="", size=12, color=TEXT, bold=False,
             font=FONT_BODY, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.03, fit=True, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = set_text_frame(box.text_frame, margin, valign, fit)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    paragraph.space_after = Pt(0)
    style_run(paragraph.add_run(), size, color, bold, font).text = text
    return box


def card_text(slide, x, y, w, h, lines, fill="FFFFFF", line=BORDER, padding=0.14):
    shape = add_shape(slide, x, y, w, h, fill, line)
    tf = set_text_frame(shape.text_frame, padding, MSO_ANCHOR.TOP, True)
    for index, item in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_before = Pt(item.get("before", 0))
        paragraph.space_after = Pt(item.get("after", 3))
        paragraph.line_spacing = item.get("line", 1.0)
        run = paragraph.add_run()
        run.text = item.get("text", "")
        style_run(
            run,
            item.get("size", 11),
            item.get("color", TEXT),
            item.get("bold", False),
            item.get("font", FONT_BODY),
            item.get("italic", False),
            item.get("underline", False),
        )
        if item.get("hyperlink"):
            run.hyperlink.address = item["hyperlink"]
    return shape


def add_footer(slide, number, total, run_id):
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.15), Inches(0.18), Inches(0.025)).fill.solid()
    rule = slide.shapes[-1]
    rule.fill.fore_color.rgb = rgb(RED)
    rule.line.fill.background()
    text_box(slide, 0.72, 7.06, 5.8, 0.2, f"VSF • Team CX • {run_id}", 7.5, MUTED)
    text_box(slide, 12.1, 7.04, 0.75, 0.22, f"{number:02d} / {total:02d}", 8, MUTED, True, align=PP_ALIGN.RIGHT)


def add_logo(slide, logo_path: Path):
    add_shape(slide, 11.82, 0.33, 1.0, 0.54, "FFFFFF", BORDER)
    slide.shapes.add_picture(str(logo_path), Inches(12.02), Inches(0.41), width=Inches(0.6), height=Inches(0.38))


def add_header(slide, eyebrow, title, logo_path, subtitle=""):
    text_box(slide, 0.5, 0.3, 10.9, 0.2, eyebrow, 8.5, RED, True, FONT_HEAD)
    title_size = 20 if len(title) < 100 else 17.5 if len(title) < 145 else 15.5
    text_box(slide, 0.5, 0.54, 10.95, 0.64, title, title_size, TEXT, True, FONT_HEAD, fit=True, line_spacing=0.9)
    if subtitle:
        text_box(slide, 0.5, 1.13, 10.8, 0.2, subtitle, 8, MUTED)
    add_logo(slide, logo_path)


def add_picture_fit(slide, path: Path, x, y, w, h, contain=False, hyperlink=""):
    with Image.open(path) as image:
        iw, ih = image.size
    if contain:
        scale = min(w / iw, h / ih)
        pw, ph = iw * scale, ih * scale
        pic = slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))
    else:
        scale = max(w / iw, h / ih)
        pw, ph = iw * scale, ih * scale
        pic = slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))
        pic.crop_left = max(0, (pw - w) / 2 / pw)
        pic.crop_right = pic.crop_left
        pic.crop_top = max(0, (ph - h) / 2 / ph)
        pic.crop_bottom = pic.crop_top
        pic.left, pic.top, pic.width, pic.height = Inches(x), Inches(y), Inches(w), Inches(h)
    if hyperlink:
        pic.click_action.hyperlink.address = hyperlink
    return pic


def add_labeled_field(tf, label, value, size=9.5, bullets=None):
    paragraph = tf.add_paragraph()
    paragraph.space_before = Pt(3)
    paragraph.space_after = Pt(1)
    style_run(paragraph.add_run(), size, TEXT, True).text = label
    if value:
        style_run(paragraph.add_run(), size, TEXT).text = f": {value}"
    for value in bullets or []:
        bullet = tf.add_paragraph()
        bullet.level = 0
        bullet.text = f"• {value[:1].upper() + value[1:] if value else value}"
        bullet.space_after = Pt(1)
        for run in bullet.runs:
            style_run(run, size - 0.5, TEXT)


def get_image_paths(asset_dir: Path):
    return {
        path.stem.upper(): path
        for path in asset_dir.iterdir()
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"} and path.stem.upper().startswith("NEWS-")
    }


def parse_findings(gen, report, overlay, images):
    result = []
    section = gen.section_by_name(report, "Findings")
    builder = gen.HtmlDeckBuilder(report, "", report.run_id, news_images={key: {"uri": "", "width": 1, "height": 1, "layout": "top"} for key in images}, overlay=overlay)
    for block in section.blocks:
        signal_id, title = gen.split_heading(block.heading)
        leading, subsections = gen.split_h4(block.lines)
        signal = gen.field_text(gen.parse_fields(leading), "Signal")
        news = gen.parse_news(gen.subsection_by_name(subsections, "News"))
        for item in news:
            subtitle, connection, published, _ = builder.editorial_news_copy(item, signal_id)
            item.update({"subtitle": subtitle, "connection": connection, "published": published, "highlights": overlay.news[item["id"]].get("highlights", []) if overlay else []})
        result.append({
            "signal_id": signal_id,
            "title": title,
            "signal": signal,
            "news": news,
            "ots": gen.parse_bold_records(gen.subsection_by_name(subsections, "Opportunity / Threat")),
            "maps": gen.parse_bold_records(gen.subsection_by_name(subsections, "Product Mapping")),
            "gaps": gen.parse_bold_records(gen.subsection_by_name(subsections, "Product Gap")),
            "actions": gen.parse_bold_records(gen.subsection_by_name(subsections, "Action")),
        })
    return result


def build_cover(prs, report, gen, logo_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, W, H, "FFFFFF", "FFFFFF", False)
    add_shape(slide, 0, 0, W, 0.08, RED, RED, False)
    add_shape(slide, 0, 0, 0.25, H, RED, RED, False)
    add_shape(slide, 0.25, 0, 0.06, H, DEEP_RED, DEEP_RED, False)
    text_box(slide, 0.78, 2.32, 7.8, 0.25, "VSF MARKET INTELLIGENCE", 10, RED, True, FONT_HEAD)
    text_box(slide, 0.78, 2.63, 7.8, 1.05, "Market Intelligence\nReport", 38, TEXT, True, FONT_HEAD, fit=True)
    text_box(slide, 0.78, 4.12, 9.5, 0.62, "Phòng Nghiên cứu thị trường và Trải nghiệm khách hàng\n• Khối Smart City", 18, MUTED)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(4.68), Inches(1.38), Inches(0.07))
    rule.fill.solid(); rule.fill.fore_color.rgb = rgb(RED); rule.line.fill.background()
    text_box(slide, 0.78, 4.99, 7.5, 0.42, gen.presentation_cover_date(report.crawl_window), 15, TEXT, True)
    add_shape(slide, 11.48, 0.30, 1.45, 0.62, "FFFFFF", BORDER)
    slide.shapes.add_picture(str(logo_path), Inches(11.70), Inches(0.46), width=Inches(1.0), height=Inches(0.29))


def build_exec(prs, gen, report, logo_path, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "01 · EXECUTIVE SUMMARY", "Điểm nhấn thị trường và Đề xuất hành động", logo_path)
    section = gen.section_by_name(report, "Executive Summary")
    y = 1.5
    row_h = 1.55
    for block in section.blocks:
        signal_id = gen.extract_ids(block.heading, "SIGNAL")[0]
        action_id = gen.extract_ids(block.heading, "ACTION")[0]
        fields = gen.parse_fields(block.lines)
        signal = gen.finding_board_copy(gen.field_text(fields, "Signal"))
        action = gen.finding_board_copy(gen.field_text(fields, "Action"))
        response = gen.enum_prefix(gen.field_text(fields, "Hướng phản hồi"), ("PREPARE", "VALIDATE", "MONITOR", "ACT"), "response")
        priority = gen.enum_prefix(gen.field_text(fields, "Priority"), ("CRITICAL", "HIGH", "MEDIUM", "LOW"), "priority")
        card_text(slide, 0.5, y, 5.25, row_h, [
            {"text": signal_id, "size": 9, "bold": True, "color": MUTED, "font": FONT_HEAD},
            {"text": signal, "size": 11.5, "bold": False, "after": 0},
        ], SIGNAL_FILL, SIGNAL_LINE)
        text_box(slide, 5.85, y + 0.48, 0.55, 0.4, "→", 23, RED, True, FONT_HEAD, align=PP_ALIGN.CENTER)
        card_text(slide, 6.45, y, 6.38, row_h, [
            {"text": f"{action_id} · {response} · {priority}", "size": 9, "bold": True, "color": RED, "font": FONT_HEAD},
            {"text": action, "size": 11.5, "bold": False, "after": 0},
        ], SOFT, BORDER)
        y += row_h + 0.22
    add_footer(slide, 2, total, report.run_id)


def add_news_card(slide, item, image_path, x, y, w, h, with_image=False):
    shape = add_shape(slide, x, y, w, h, "FFFFFF", BORDER)
    shape.text_frame.clear()
    image_h = min(1.05, h * 0.24) if with_image and image_path else 0
    copy_h = h - 0.42 - (image_h + 0.18 if image_h else 0)
    copy = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.1), Inches(w - 0.24), Inches(copy_h))
    tf = set_text_frame(copy.text_frame, 0.01, MSO_ANCHOR.TOP, True)
    title = tf.paragraphs[0]
    style_run(title.add_run(), 10.2 if h > 3 else 8.7, TEXT, True, FONT_HEAD).text = item["subtitle"]
    connection = tf.add_paragraph(); connection.space_before = Pt(3); connection.space_after = Pt(2)
    style_run(connection.add_run(), 8.4, RED, True).text = f"Liên hệ {item.get('signal_id', '')}: "
    style_run(connection.add_run(), 8.4, TEXT, True).text = item["connection"]
    summary = tf.add_paragraph(); summary.space_before = Pt(2); summary.space_after = Pt(2)
    raw = item["summary"]
    highlights = list(item.get("highlights", []))
    spans = []
    for phrase in highlights:
        start = raw.find(phrase)
        if start >= 0:
            spans.append((start, start + len(phrase)))
    spans.sort()
    cursor = 0
    for start, end in spans:
        style_run(summary.add_run(), 8.7 if h > 3 else 7.7, TEXT).text = raw[cursor:start]
        style_run(summary.add_run(), 8.7 if h > 3 else 7.7, TEXT, True, underline=True).text = raw[start:end]
        cursor = end
    style_run(summary.add_run(), 8.7 if h > 3 else 7.7, TEXT).text = raw[cursor:]
    if image_h:
        add_picture_fit(slide, image_path, x + 0.13, y + h - image_h - 0.4, w - 0.26, image_h, contain=False, hyperlink=item.get("source_url", ""))
    footer = text_box(slide, x + 0.13, y + h - 0.28, w - 0.26, 0.17, "", 7.4, MUTED, margin=0, fit=True).text_frame.paragraphs[0]
    run = style_run(footer.add_run(), 7.4, MUTED)
    run.text = f"{item['id']} · {item['source_name']} ↗   Xuất bản: {item['published']}"
    run.hyperlink.address = item.get("source_url", "")
    return shape


def build_page_a(prs, gen, finding, images, logo_path, slide_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = f"{finding['signal_id']} — {gen.finding_board_copy(finding['title'])}"
    index = slide_no // 2
    add_header(slide, f"{index:02d} · FINDING {index:02d}A", title, logo_path, "Signal, Evidence và Opportunity / Threat")
    signal = add_shape(slide, 0.5, 1.38, 12.32, 0.58, SIGNAL_FILL, SIGNAL_LINE)
    tf = set_text_frame(signal.text_frame, 0.16, MSO_ANCHOR.MIDDLE, True)
    p = tf.paragraphs[0]
    style_run(p.add_run(), 9, RED, True, FONT_HEAD).text = "SIGNAL   "
    style_run(p.add_run(), 10.5, BLUE_TEXT).text = gen.finding_board_copy(finding["signal"])
    text_box(slide, 0.5, 2.05, 7.75, 0.2, "EVIDENCE", 8.5, MUTED, True, FONT_HEAD)
    news = finding["news"]
    singleton = {0} if len(news) in {1, 3} else set(range(len(news))) if len(news) == 2 else set()
    for item in news:
        item["signal_id"] = finding["signal_id"]
    if len(news) == 1:
        positions = [(0.5, 2.32, 7.75, 4.55)]
    elif len(news) == 2:
        positions = [(0.5, 2.32, 3.78, 4.55), (4.47, 2.32, 3.78, 4.55)]
    elif len(news) == 3:
        positions = [(0.5, 2.32, 3.78, 4.55), (4.47, 2.32, 3.78, 2.17), (4.47, 4.7, 3.78, 2.17)]
    else:
        positions = [(0.5, 2.32, 3.78, 2.17), (4.47, 2.32, 3.78, 2.17), (0.5, 4.7, 3.78, 2.17), (4.47, 4.7, 3.78, 2.17)]
    inline_ids = set()
    for idx, (item, pos) in enumerate(zip(news, positions)):
        with_image = idx in singleton and item["id"] in images
        if with_image:
            inline_ids.add(item["id"])
        add_news_card(slide, item, images.get(item["id"]), *pos, with_image)
    gallery = [item for item in news if item["id"] in images and item["id"] not in inline_ids]
    right_x, right_w = 8.45, 4.37
    ot_y = 2.32
    if gallery:
        gallery_h = 1.65
        frame = add_shape(slide, right_x, 2.32, right_w, gallery_h, SOFT, BORDER)
        del frame
        each_w = right_w / min(len(gallery), 2)
        for idx, item in enumerate(gallery[:2]):
            add_picture_fit(slide, images[item["id"]], right_x + idx * each_w + 0.05, 2.37, each_w - 0.1, gallery_h - 0.1, contain=False, hyperlink=item.get("source_url", ""))
        ot_y = 4.13
    ot_records = finding["ots"]
    available = 6.87 - ot_y
    ot_h = (available - 0.12 * (len(ot_records) - 1)) / len(ot_records)
    for index, (heading, fields) in enumerate(ot_records):
        ot_id, kind = gen.record_heading(heading)
        priority = gen.clean_atom(gen.field_text(fields, "Mức độ quan trọng"))
        threat = gen.norm(kind) == "threat"
        card_text(slide, right_x, ot_y + index * (ot_h + 0.12), right_w, ot_h, [
            {"text": f"{ot_id} {kind.upper()} · {priority}", "size": 8.5, "bold": True, "color": DEEP_RED if threat else "075E54", "font": FONT_HEAD},
            {"text": gen.finding_board_copy(gen.field_text(fields, "Nội dung")), "size": 9.4, "color": DEEP_RED if threat else "075E54", "after": 0},
        ], THREAT_FILL if threat else OPP_FILL, THREAT_LINE if threat else OPP_LINE)
    add_footer(slide, slide_no, total, finding.get("run_id", ""))


def build_page_b(prs, gen, finding, logo_path, slide_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = f"{finding['signal_id']} — {gen.finding_board_copy(finding['title'])}"
    index = (slide_no - 1) // 2
    add_header(slide, f"{index:02d} · FINDING {index:02d}B", title, logo_path, "Nhu cầu thị trường, Khoảng trống sản phẩm và Đề xuất hành động đã duyệt")
    map_heading, map_fields = finding["maps"][0]
    gap_heading, gap_fields = finding["gaps"][0]
    action_heading, action_fields = finding["actions"][0]
    pm_id, pm_title = gen.record_heading(map_heading)
    gap_id, _ = gen.record_heading(gap_heading)
    action_id, priority = gen.record_heading(action_heading)
    response = gen.clean_atom(gen.field_text(action_fields, "Hướng phản hồi")).upper()
    map_card = add_shape(slide, 0.5, 1.45, 5.55, 4.05, "FFFFFF", BORDER)
    tf = set_text_frame(map_card.text_frame, 0.18, MSO_ANCHOR.TOP, True)
    p = tf.paragraphs[0]; style_run(p.add_run(), 10, RED, True, FONT_HEAD).text = f"NHU CẦU THỊ TRƯỜNG                         {pm_id}"
    p = tf.add_paragraph(); p.space_before = Pt(4); style_run(p.add_run(), 13, TEXT, True, FONT_HEAD).text = gen.finding_board_copy(pm_title)
    add_labeled_field(tf, "O/T liên quan", gen.field_text(map_fields, "O/T liên quan"))
    add_labeled_field(tf, "Vấn đề thị trường", gen.finding_board_copy(gen.field_text(map_fields, "Vấn đề thị trường")))
    mandatory = gen.field_bullets(map_fields, "Năng lực bắt buộc")
    if mandatory:
        add_labeled_field(tf, "Năng lực bắt buộc", "", 9.2, [gen.finding_board_copy(value) for value in mandatory])
    add_labeled_field(tf, "Khách hàng mục tiêu", gen.finding_board_copy(gen.field_text(map_fields, "Khách hàng mục tiêu")))
    gap_card = add_shape(slide, 6.25, 1.45, 6.57, 4.05, "FFFFFF", BORDER)
    tf = set_text_frame(gap_card.text_frame, 0.18, MSO_ANCHOR.TOP, True)
    status = gen.clean_atom(gen.field_text(gap_fields, "Trạng thái capability"))
    level = gen.clean_atom(gen.field_text(gap_fields, "Mức độ gap"))
    p = tf.paragraphs[0]; style_run(p.add_run(), 10, RED, True, FONT_HEAD).text = f"KHOẢNG TRỐNG SẢN PHẨM                 {gap_id} · {status} · {level}"
    add_labeled_field(tf, "Product Mapping liên quan", gen.field_text(gap_fields, "Product Mapping liên quan"), 9.2)
    add_labeled_field(tf, "Sản phẩm VSF liên quan", gen.finding_board_copy(gen.field_text(gap_fields, "Sản phẩm VSF liên quan")), 9.2)
    missing = [gen.finding_board_copy(value) for value in gen.split_semicolon_items([gen.field_text(gap_fields, "Capability còn thiếu")])]
    add_labeled_field(tf, "Tính năng còn thiếu", "", 9.0, missing)
    action = add_shape(slide, 0.5, 5.68, 12.32, 1.17, SOFT, BORDER)
    tf = set_text_frame(action.text_frame, 0.18, MSO_ANCHOR.TOP, True)
    p = tf.paragraphs[0]; style_run(p.add_run(), 9, RED, True, FONT_HEAD).text = f"{action_id} · {response} · {priority}     Product Gap: {gen.field_text(action_fields, 'Product Gap liên quan')}"
    values = [
        ("Đề xuất", gen.finding_board_copy(gen.field_text(action_fields, "Hành động đề xuất"))),
        ("Bước tiếp theo", gen.finding_board_copy(gen.field_text(action_fields, "Bước tiếp theo"))),
        ("Kết quả mong đợi", gen.finding_board_copy(gen.field_text(action_fields, "Kết quả mong đợi"))),
    ]
    p = tf.add_paragraph(); p.space_before = Pt(3)
    for idx, (label, value) in enumerate(values):
        if idx:
            style_run(p.add_run(), 8.2, MUTED).text = "     |     "
        style_run(p.add_run(), 8.4, TEXT, True).text = f"{label}: "
        style_run(p.add_run(), 8.4, TEXT).text = value
    add_footer(slide, slide_no, total, finding.get("run_id", ""))


def build_approach(prs, gen, report, logo_path, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "03 · APPROACH", "Cách chuyển Signal thành Action", logo_path, "Giữ đầy đủ lineage và chỉ sử dụng nội dung đã qua review bắt buộc")
    section = gen.section_by_name(report, "Approach")
    flow_block = next(block for block in section.blocks if gen.norm(block.heading) == gen.norm("Từ Signal đến Action"))
    response_block = next(block for block in section.blocks if gen.norm(block.heading) == gen.norm("Cách đọc hướng phản hồi"))
    text_box(slide, 0.5, 1.52, 7.0, 0.25, "TỪ SIGNAL ĐẾN ACTION", 9, MUTED, True, FONT_HEAD)
    y = 1.88
    for index, entry in enumerate(gen.parse_fields(flow_block.lines), 1):
        card_text(slide, 0.5, y, 7.0, 0.82, [
            {"text": f"{index:02d}  {entry['label']}", "size": 9, "bold": True, "color": RED, "font": FONT_HEAD},
            {"text": str(entry["value"]), "size": 9.1, "after": 0},
        ], "FFFFFF", BORDER)
        y += 0.94
    text_box(slide, 7.8, 1.52, 5.02, 0.25, "CÁCH ĐỌC HƯỚNG PHẢN HỒI", 9, MUTED, True, FONT_HEAD)
    y = 1.88
    for entry in gen.parse_fields(response_block.lines):
        response = str(entry["label"]).upper()
        color = DEEP_RED if response == "VALIDATE" else MUTED if response == "MONITOR" else RED
        card_text(slide, 7.8, y, 5.02, 1.05, [
            {"text": response, "size": 9, "bold": True, "color": color, "font": FONT_HEAD},
            {"text": str(entry["value"]), "size": 9.2, "after": 0},
        ], SOFT, BORDER)
        y += 1.18
    add_footer(slide, total, total, report.run_id)


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--source-deck", type=Path)
    args = parser.parse_args()
    script_dir = Path(__file__).resolve().parent
    gen = load_html_generator(script_dir / "generate_vsf_mi_finding_slides.py")
    input_path = args.input.resolve()
    run_root = gen.find_run_root(input_path)
    run_id = gen.validate_gates(run_root)
    report = gen.parse_report(input_path)
    if report.run_id and report.run_id != run_id:
        raise ValueError(f"Run ID mismatch: {report.run_id} != {run_id}")
    source_deck = args.source_deck.resolve() if args.source_deck else None
    overlay = gen.load_presentation_overlay(source_deck) if source_deck else None
    slide_dir = run_root / "deliverables" / "slides"
    images = get_image_paths(slide_dir / "assets")
    logo_path = script_dir.parent / "assets" / "vsf-logo-transparent.png"
    findings = parse_findings(gen, report, overlay, images)
    total = 3 + 2 * len(findings)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    build_cover(prs, report, gen, logo_path)
    build_exec(prs, gen, report, logo_path, total)
    slide_no = 3
    for finding in findings:
        finding["run_id"] = run_id
        build_page_a(prs, gen, finding, images, logo_path, slide_no, total)
        slide_no += 1
        build_page_b(prs, gen, finding, logo_path, slide_no, total)
        slide_no += 1
    build_approach(prs, gen, report, logo_path, total)
    if len(prs.slides) != total:
        raise ValueError(f"Expected {total} slides, created {len(prs.slides)}")
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print({
        "output": str(args.output.resolve()),
        "slides": len(prs.slides),
        "editable_text_shapes": sum(1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)),
        "images": sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13),
        "gate_check": "APPROVED",
        "run_id": run_id,
    })


if __name__ == "__main__":
    main()
