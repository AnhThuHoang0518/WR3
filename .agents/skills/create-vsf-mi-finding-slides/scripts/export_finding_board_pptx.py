"""Package a rendered 16:9 Finding Board PDF as a visually faithful PPTX."""

from __future__ import annotations

import argparse
import io
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Inches


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", required=True, type=Path, help="Rendered Finding Board PDF")
    parser.add_argument("--output", required=True, type=Path, help="PPTX output path")
    parser.add_argument("--dpi", type=int, default=192, help="Raster resolution per PDF page")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if not args.input.is_file():
        raise FileNotFoundError(args.input)
    if args.dpi < 96:
        raise ValueError("--dpi must be at least 96")

    document = fitz.open(args.input)
    if document.page_count == 0:
        raise ValueError("The input PDF has no pages")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    scale = args.dpi / 72

    for page in document:
        image = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(image.tobytes("png")),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    if len(presentation.slides) != document.page_count:
        raise RuntimeError("PPTX slide count does not match PDF page count")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(args.output)
    print(
        {
            "output": str(args.output.resolve()),
            "slides": len(presentation.slides),
            "source_pages": document.page_count,
            "dpi": args.dpi,
            "aspect_ratio": "16:9",
        }
    )


if __name__ == "__main__":
    main()
